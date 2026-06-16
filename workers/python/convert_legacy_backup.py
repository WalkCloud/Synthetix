"""Synthetix document converter — converts files to Markdown with image extraction.

Usage: python convert.py <input_file> <output_dir>
Output: writes full.md to output_dir (with image references), prints output path to stdout
        extracts images to <output_dir>/images/
"""
import sys
import os
import base64
import hashlib
import re
from pathlib import Path

from markitdown import MarkItDown


def _ensure_images_dir(output_dir: str) -> str:
    images_dir = os.path.join(output_dir, "images")
    os.makedirs(images_dir, exist_ok=True)
    return images_dir


def _save_image(image_bytes: bytes, images_dir: str, index: int, ext: str = "png") -> str:
    """Save image bytes to disk, return relative markdown path."""
    h = hashlib.md5(image_bytes).hexdigest()[:8]
    filename = f"img_{index:03d}_{h}.{ext}"
    filepath = os.path.join(images_dir, filename)
    with open(filepath, "wb") as f:
        f.write(image_bytes)
    return f"images/{filename}"


def _image_ext_from_content_type(content_type: str) -> str:
    mapping = {
        "image/png": "png",
        "image/jpeg": "jpg",
        "image/jpg": "jpg",
        "image/gif": "gif",
        "image/webp": "webp",
        "image/svg+xml": "svg",
        "image/tiff": "tiff",
        "image/bmp": "bmp",
    }
    return mapping.get(content_type.lower(), "png")


def _image_ext_from_name(name: str) -> str:
    ext = os.path.splitext(name)[1].lower().lstrip(".")
    return ext if ext in ("png", "jpg", "jpeg", "gif", "webp", "svg", "tiff", "bmp") else "png"


def convert_docx(input_file: str, output_dir: str) -> str:
    """Convert DOCX to markdown with embedded images extracted."""
    from docx import Document as DocxDocument

    images_dir = _ensure_images_dir(output_dir)
    doc = DocxDocument(input_file)
    md_lines = []
    img_index = 0

    # Build a set of image relationship IDs for quick lookup
    image_rels = {}
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image_rels[rel.rId] = rel

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            # Paragraph
            para = element
            from docx.oxml.ns import qn
            text_parts = []
            has_image = False

            for child in para:
                child_tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

                if child_tag == "r":
                    # Run
                    for run_child in child:
                        run_tag = run_child.tag.split("}")[-1] if "}" in run_child.tag else run_child.tag
                        if run_tag == "t":
                            text_parts.append(run_child.text or "")
                        elif run_tag == "drawing" or run_tag == "pict":
                            # Image in this run
                            has_image = True
                            # Extract image reference from drawing
                            blips = run_child.findall(f".//{qn('a:blip')}")
                            for blip in blips:
                                embed = blip.get(qn("r:embed"))
                                if embed and embed in image_rels:
                                    rel = image_rels[embed]
                                    try:
                                        img_bytes = rel.target_part.blob
                                        ext = _image_ext_from_name(rel.target_ref)
                                        rel_path = _save_image(img_bytes, images_dir, img_index, ext)
                                        md_lines.append(f"\n![Image {img_index + 1}]({rel_path})\n")
                                        img_index += 1
                                    except Exception as e:
                                        print(f"Warning: DOCX image extraction failed: {e}", file=sys.stderr)

            text = "".join(text_parts).strip()
            if text:
                # Detect heading style
                from docx.oxml.ns import qn as qn2
                pPr = element.find(qn2("w:pPr"))
                if pPr is not None:
                    pStyle = pPr.find(qn2("w:pStyle"))
                    if pStyle is not None:
                        style_val = pStyle.get(qn2("w:val"), "")
                        if style_val.startswith("Heading"):
                            level = style_val.replace("Heading", "")
                            try:
                                level = int(level)
                            except ValueError:
                                level = 1
                            md_lines.append(f"{'#' * level} {text}\n")
                            continue
                md_lines.append(text + "\n")

        elif tag == "tbl":
            # Table — extract as markdown table
            from docx.oxml.ns import qn as qn3
            rows = element.findall(f".//{qn3('w:tr')}")
            table_data = []
            for row in rows:
                cells = row.findall(f".//{qn3('w:tc')}")
                row_data = []
                for cell in cells:
                    cell_text = "".join(t.text or "" for t in cell.iter(qn3("w:t")))
                    row_data.append(cell_text.strip())
                table_data.append(row_data)

            if table_data:
                # Header row
                header = table_data[0]
                md_lines.append("| " + " | ".join(header) + " |")
                md_lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                for row in table_data[1:]:
                    # Pad row to match header length
                    while len(row) < len(header):
                        row.append("")
                    md_lines.append("| " + " | ".join(row[:len(header)]) + " |")
                md_lines.append("")

    markdown = "\n".join(md_lines).strip()

    if not markdown:
        # Fallback to MarkItDown
        return convert_generic(input_file, output_dir)

    output_path = os.path.join(output_dir, "full.md")
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(markdown)

    return output_path


def convert_pdf(input_file: str, output_dir: str) -> str:
    """Convert PDF to markdown with images extracted per page."""
    import fitz  # PyMuPDF

    images_dir = _ensure_images_dir(output_dir)
    doc = fitz.open(input_file)
    md_parts = []
    img_index = 0

    for page_num in range(len(doc)):
        page = doc[page_num]

        # Extract text
        text = page.get_text("text")
        if text.strip():
            md_parts.append(text.strip())

        # Extract images from this page
        image_list = page.get_images(full=True)
        for img_info in image_list:
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
                if base_image:
                    img_bytes = base_image["image"]
                    ext = base_image.get("ext", "png")
                    if ext == "jpeg":
                        ext = "jpg"
                    # Skip tiny images (likely decorative)
                    if len(img_bytes) < 500:
                        continue
                    rel_path = _save_image(img_bytes, images_dir, img_index, ext)
                    md_parts.append(f"\n![Page {page_num + 1} - Image {img_index + 1}]({rel_path})\n")
                    img_index += 1
            except Exception as e:
                print(f"Warning: PDF image extraction failed: {e}", file=sys.stderr)
                continue

    doc.close()

    markdown = "\n\n".join(md_parts)

    if not markdown:
        return convert_generic(input_file, output_dir)

    output_path = os.path.join(output_dir, "full.md")
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(markdown)

    return output_path


def convert_pptx(input_file: str, output_dir: str) -> str:
    """Convert PPTX to markdown with images extracted per slide."""
    from pptx import Presentation
    from pptx.util import Inches

    images_dir = _ensure_images_dir(output_dir)
    prs = Presentation(input_file)
    md_parts = []
    img_index = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)

            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    img_bytes = image.blob
                    ext = _image_ext_from_content_type(image.content_type)
                    rel_path = _save_image(img_bytes, images_dir, img_index, ext)
                    slide_texts.append(f"\n![Slide {slide_num} - Image {img_index + 1}]({rel_path})\n")
                    img_index += 1
                except Exception as e:
                    print(f"Warning: PPTX image extraction failed: {e}", file=sys.stderr)

        if slide_texts:
            md_parts.append(f"## Slide {slide_num}\n\n" + "\n\n".join(slide_texts))

    markdown = "\n\n".join(md_parts)

    if not markdown:
        return convert_generic(input_file, output_dir)

    output_path = os.path.join(output_dir, "full.md")
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(markdown)

    return output_path


def convert_xlsx(input_file: str, output_dir: str) -> str:
    """Convert XLSX to markdown with images extracted."""
    # For XLSX, use MarkItDown for text and try to extract charts/images
    # XLSX images are complex (charts, embedded objects), fall back to MarkItDown
    return convert_generic(input_file, output_dir)


def convert_html(input_file: str, output_dir: str) -> str:
    """Convert HTML to markdown. Images are referenced by URL, not extracted."""
    return convert_generic(input_file, output_dir)


def convert_epub(input_file: str, output_dir: str) -> str:
    """Convert EPUB to markdown. Images may be embedded in the EPUB archive."""
    import zipfile

    images_dir = _ensure_images_dir(output_dir)

    # EPUB is a ZIP archive — extract images from it
    img_index = 0
    try:
        with zipfile.ZipFile(input_file, "r") as zf:
            for name in zf.namelist():
                ext = os.path.splitext(name)[1].lower()
                if ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"):
                    try:
                        img_bytes = zf.read(name)
                        if len(img_bytes) < 200:
                            continue
                        ext_clean = ext.lstrip(".")
                        if ext_clean == "jpeg":
                            ext_clean = "jpg"
                        rel_path = _save_image(img_bytes, images_dir, img_index, ext_clean)
                        img_index += 1
                    except Exception:
                        continue
    except zipfile.BadZipFile:
        pass

    # Use MarkItDown for text content
    return convert_generic(input_file, output_dir)


def convert_generic(input_file: str, output_dir: str) -> str:
    """Fallback: use MarkItDown for text-only conversion."""
    md = MarkItDown()
    result = md.convert(input_file)
    output_path = os.path.join(output_dir, "full.md")

    with open(output_path, "w", encoding="utf-8") as f:
        f.write(result.text_content)

    return output_path


def main():
    if len(sys.argv) != 3:
        print("Usage: python convert.py <input_file> <output_dir>", file=sys.stderr)
        sys.exit(1)

    input_file = sys.argv[1]
    output_dir = sys.argv[2]

    if not os.path.exists(input_file):
        print(f"Input file not found: {input_file}", file=sys.stderr)
        sys.exit(1)

    os.makedirs(output_dir, exist_ok=True)

    ext = os.path.splitext(input_file)[1].lower()

    converters = {
        ".docx": convert_docx,
        ".pdf": convert_pdf,
        ".pptx": convert_pptx,
        ".xlsx": convert_xlsx,
        ".html": convert_html,
        ".htm": convert_html,
        ".epub": convert_epub,
    }

    converter = converters.get(ext, convert_generic)
    output_path = converter(input_file, output_dir)

    print(output_path)


if __name__ == "__main__":
    main()
